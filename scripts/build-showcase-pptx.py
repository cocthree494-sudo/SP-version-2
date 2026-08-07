"""Build the Relay showcase PPTX from browser-rendered 16:9 slide images."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


SLIDE_COUNT = 27


def build_deck(slide_dir: Path, output: Path) -> None:
    images = [slide_dir / f"slide-{index:02d}.png" for index in range(1, SLIDE_COUNT + 1)]
    missing = [image for image in images if not image.is_file()]
    if missing:
        raise FileNotFoundError(f"Missing rendered slide: {missing[0]}")

    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    blank_layout = deck.slide_layouts[6]

    # Remove the starter slide collection, then add one lossless image per slide.
    while deck.slides:
        slide_id = deck.slides._sldIdLst[0]
        deck.part.drop_rel(slide_id.rId)
        del deck.slides._sldIdLst[0]

    for image in images:
        slide = deck.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=deck.slide_width,
            height=deck.slide_height,
        )

    deck.core_properties.title = "Relay — Universal Support Agent Team Showcase"
    deck.core_properties.subject = "Product, architecture, security, delivery, and roadmap"
    deck.core_properties.author = "Relay team"
    deck.core_properties.keywords = "Relay, AI support, SaaS, architecture, presentation"
    output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(output)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("slide_dir", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    build_deck(args.slide_dir.resolve(), args.output.resolve())


if __name__ == "__main__":
    main()
